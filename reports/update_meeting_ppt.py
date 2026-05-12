# -*- coding: utf-8 -*-
"""기존 PPT 수정: 샘플 데이터 슬라이드에 폴더 구조 추가 + 페이지 번호 보정."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import re

PPT = r"c:\projects\iCreat\reports\2026-05-13_meeting_update.pptx"

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xE8, 0x74, 0x22)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
SUB_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE0, 0xE6, 0xF0)
TREE_GRAY = RGBColor(0x8A, 0x95, 0xA8)

FONT_KR = "맑은 고딕"
FONT_MONO = "Consolas"


def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             font=FONT_KR, size=18, bold=False, color=GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


prs = Presentation(PPT)
total = len(prs.slides)
print(f"총 슬라이드 수: {total}")

# ---------- 1) 페이지 번호 보정: "X / 13" → "X / N" ----------
page_pat = re.compile(r"^\s*(\d+)\s*/\s*\d+\s*$")
fixed = 0
for idx, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                m = page_pat.match(run.text)
                if m:
                    run.text = f"{idx} / {total}"
                    fixed += 1
print(f"페이지 번호 보정: {fixed}건")

# ---------- 2) 샘플 데이터 슬라이드 찾기 ----------
target_idx = None
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and "샘플 데이터 준비 현황" in shape.text_frame.text:
            target_idx = idx
            break
    if target_idx is not None:
        break
print(f"샘플 데이터 슬라이드 위치: {target_idx + 1 if target_idx is not None else 'NOT FOUND'}")

slide = prs.slides[target_idx]

# 기존 도형 모두 제거 (헤더부터 다시 그림)
sp_tree = slide.shapes._spTree
for sp in list(sp_tree):
    if sp.tag.endswith("}sp") or sp.tag.endswith("}pic") or sp.tag.endswith("}graphicFrame") or sp.tag.endswith("}grpSp") or sp.tag.endswith("}cxnSp"):
        sp_tree.remove(sp)

SLIDE_H = prs.slide_height

# ---------- 헤더 재구성 ----------
add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
add_text(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.6),
         "샘플 데이터 준비 현황", size=28, bold=True, color=NAVY)
add_text(slide, Inches(0.72), Inches(0.95), Inches(11), Inches(0.35),
         "Sensor_monitor/sensor_data.zip — 폴더 구조 한눈에 보기",
         size=14, color=SUB_GRAY)
add_rect(slide, Inches(0.7), Inches(1.4), Inches(0.8), Inches(0.06), ACCENT)
add_text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
         "Sensor Monitor → 베데스다 iCReaT DCT 서버 전환", size=10, color=SUB_GRAY)
add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
         f"{target_idx + 1} / {total}", size=10, color=SUB_GRAY, align=PP_ALIGN.RIGHT)

# ---------- 상단 통계 4개 (컴팩트, 가로 한 줄) ----------
stats = [
    ("175",  "전체 엔트리",   NAVY),
    ("172",  "CSV 파일",      BLUE),
    ("143",  "sended 병합본", ACCENT),
    ("3",    "log 파일",      SUB_GRAY),
]
sx = Inches(0.7); sy = Inches(1.75)
sw = Inches(2.95); sh = Inches(1.05)
for i, (n, label, color) in enumerate(stats):
    x = sx + (sw + Inches(0.1)) * i
    add_rect(slide, x, sy, sw, sh, LIGHT_BG)
    add_rect(slide, x, sy, Inches(0.1), sh, color)
    add_text(slide, x + Inches(0.3), sy + Inches(0.12), Inches(1.2), Inches(0.8),
             n, size=32, bold=True, color=color)
    add_text(slide, x + Inches(1.55), sy + Inches(0.35), Inches(1.3), Inches(0.5),
             label, size=13, color=GRAY_TEXT)

# ---------- 좌측: 폴더 구조 트리 ----------
TREE_X = Inches(0.7); TREE_Y = Inches(3.0)
TREE_W = Inches(7.6); TREE_H = Inches(3.95)
add_rect(slide, TREE_X, TREE_Y, TREE_W, TREE_H, CODE_BG)

# 트리 헤더
add_text(slide, TREE_X + Inches(0.3), TREE_Y + Inches(0.18),
         TREE_W - Inches(0.6), Inches(0.4),
         "📁 폴더 구조", size=13, bold=True, color=ACCENT, font=FONT_MONO)

# 트리 라인들 — (들여쓰기, 텍스트, 강조여부)
# 트리 문자는 mono font 로 렌더, sended 와 30분 병합본은 ACCENT 색
tree_lines = [
    ("sensor_data/",                                 "dir_root"),
    ("├─ logs/",                                     "dir"),
    ("│   ├─ app_debug_log_260327.txt",              "file"),
    ("│   ├─ app_debug_log_260328.txt",              "file"),
    ("│   └─ app_debug_log_260330.txt",              "file"),
    ("├─ HeartRate/                  ← 센서 폴더",   "dir_sensor"),
    ("│   ├─ HeartRate_260330_190548.csv     ← 5분 조각",  "file"),
    ("│   ├─ HeartRate_260330_191048.csv",           "file"),
    ("│   ├─  …",                                    "file_dim"),
    ("│   └─ sended/   ← 전송 성공한 30분 병합본",   "dir_sended"),
    ("│       ├─ HeartRate_260330_1130.csv",         "file_sended"),
    ("│       ├─ HeartRate_260330_1200.csv",         "file_sended"),
    ("│       └─  …",                                "file_sended_dim"),
    ("├─ Accelerometer/  …  (동일 구조: 조각 + sended/)",  "dir_etc"),
    ("├─ Gyroscope/      …",                         "dir_etc"),
    ("├─ Gravity/        …",                         "dir_etc"),
    ("├─ Light/          …",                         "dir_etc"),
    ("├─ PpgGreen/       …",                         "dir_etc"),
    ("└─ StepCount/      …",                         "dir_etc"),
]

tb = slide.shapes.add_textbox(
    TREE_X + Inches(0.3), TREE_Y + Inches(0.6),
    TREE_W - Inches(0.6), TREE_H - Inches(0.75))
tf = tb.text_frame
tf.word_wrap = False
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

for i, (line, kind) in enumerate(tree_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1)

    if kind == "dir_sended" or kind == "file_sended" or kind == "file_sended_dim":
        # 트리 가지(│ ├ └)는 회색, 파일/디렉토리명은 ACCENT
        tree_part = ""
        name_part = line
        for ch in line:
            if ch in "│├└─ ":
                tree_part += ch
            else:
                break
        name_part = line[len(tree_part):]
        if tree_part:
            r1 = p.add_run(); r1.text = tree_part
            r1.font.name = FONT_MONO; r1.font.size = Pt(13)
            r1.font.color.rgb = TREE_GRAY
        r2 = p.add_run(); r2.text = name_part
        r2.font.name = FONT_MONO; r2.font.size = Pt(13)
        r2.font.color.rgb = ACCENT
        if kind != "file_sended_dim":
            r2.font.bold = True
    elif kind == "dir_root":
        r = p.add_run(); r.text = line
        r.font.name = FONT_MONO; r.font.size = Pt(14)
        r.font.color.rgb = WHITE; r.font.bold = True
    elif kind == "dir_sensor":
        # 트리 + 이름 + 주석 분리
        tree_part = ""
        for ch in line:
            if ch in "│├└─ ":
                tree_part += ch
            else:
                break
        rest = line[len(tree_part):]
        if "←" in rest:
            name_part, comment_part = rest.split("←", 1)
            r1 = p.add_run(); r1.text = tree_part
            r1.font.name = FONT_MONO; r1.font.size = Pt(13)
            r1.font.color.rgb = TREE_GRAY
            r2 = p.add_run(); r2.text = name_part
            r2.font.name = FONT_MONO; r2.font.size = Pt(13)
            r2.font.color.rgb = WHITE; r2.font.bold = True
            r3 = p.add_run(); r3.text = "←" + comment_part
            r3.font.name = FONT_MONO; r3.font.size = Pt(11)
            r3.font.color.rgb = TREE_GRAY
        else:
            r = p.add_run(); r.text = line
            r.font.name = FONT_MONO; r.font.size = Pt(13)
            r.font.color.rgb = WHITE
    elif kind == "dir_etc" or kind == "dir":
        r = p.add_run(); r.text = line
        r.font.name = FONT_MONO; r.font.size = Pt(13)
        r.font.color.rgb = CODE_FG
    elif kind == "file_dim":
        r = p.add_run(); r.text = line
        r.font.name = FONT_MONO; r.font.size = Pt(13)
        r.font.color.rgb = TREE_GRAY
    elif kind == "file":
        # 5분 조각 — 트리는 회색, 파일명은 옅은 색, 주석은 회색
        tree_part = ""
        for ch in line:
            if ch in "│├└─ ":
                tree_part += ch
            else:
                break
        rest = line[len(tree_part):]
        if "←" in rest:
            name_part, comment_part = rest.split("←", 1)
            r1 = p.add_run(); r1.text = tree_part
            r1.font.name = FONT_MONO; r1.font.size = Pt(13)
            r1.font.color.rgb = TREE_GRAY
            r2 = p.add_run(); r2.text = name_part
            r2.font.name = FONT_MONO; r2.font.size = Pt(13)
            r2.font.color.rgb = CODE_FG
            r3 = p.add_run(); r3.text = "←" + comment_part
            r3.font.name = FONT_MONO; r3.font.size = Pt(11)
            r3.font.color.rgb = TREE_GRAY
        else:
            r1 = p.add_run(); r1.text = tree_part
            r1.font.name = FONT_MONO; r1.font.size = Pt(13)
            r1.font.color.rgb = TREE_GRAY
            r2 = p.add_run(); r2.text = rest
            r2.font.name = FONT_MONO; r2.font.size = Pt(13)
            r2.font.color.rgb = CODE_FG

# ---------- 우측: 핵심 포인트 + 센서 목록 ----------
RX = Inches(8.45); RW = Inches(4.2)

# 핵심 포인트
add_rect(slide, RX, Inches(3.0), RW, Inches(2.3), LIGHT_BG)
add_rect(slide, RX, Inches(3.0), RW, Inches(0.08), ACCENT)
add_text(slide, RX + Inches(0.25), Inches(3.2), RW - Inches(0.5), Inches(0.45),
         "핵심 포인트", size=15, bold=True, color=NAVY)

points = [
    ("센서별 폴더", "7종 동일 구조"),
    ("5분 조각 CSV", "센서 폴더 바로 아래"),
    ("sended/ 하위", "30분 단위 병합본"),
    ("sended = ", "전송 성공 보관 위치"),
]
py = Inches(3.7)
for label, desc in points:
    add_text(slide, RX + Inches(0.25), py, Inches(0.2), Inches(0.3),
             "▸", size=12, bold=True, color=ACCENT)
    add_text(slide, RX + Inches(0.5), py - Inches(0.02),
             RW - Inches(0.75), Inches(0.3),
             label, size=11, bold=True, color=NAVY)
    add_text(slide, RX + Inches(0.5), py + Inches(0.22),
             RW - Inches(0.75), Inches(0.3),
             desc, size=10, color=GRAY_TEXT)
    py += Inches(0.4)

# 센서 7종 (NAVY 박스)
add_rect(slide, RX, Inches(5.45), RW, Inches(1.5), NAVY)
add_text(slide, RX + Inches(0.25), Inches(5.6), RW - Inches(0.5), Inches(0.35),
         "포함 센서 (7종)", size=13, bold=True, color=WHITE)
sensors = ["PpgGreen", "HeartRate", "Accelerometer",
           "Gyroscope", "Gravity", "Light", "StepCount"]
ssx = RX + Inches(0.25); ssy = Inches(6.0)
for i, name in enumerate(sensors):
    col = i % 2; row = i // 2
    x = ssx + Inches(1.95) * col
    y = ssy + Inches(0.28) * row
    add_text(slide, x, y, Inches(1.9), Inches(0.3),
             f"▸ {name}", size=10, color=WHITE, font=FONT_MONO)

prs.save(PPT)
print(f"saved: {PPT}")
