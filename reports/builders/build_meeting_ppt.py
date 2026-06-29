# -*- coding: utf-8 -*-
"""2026-05-13 과제 미팅 보고 PPT 생성 스크립트."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------- 색상 / 폰트 ----------------
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xE8, 0x74, 0x22)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
SUB_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE0, 0xE6, 0xF0)

FONT_KR = "맑은 고딕"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = fill_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             font=FONT_KR, size=18, bold=False, color=GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None, page_no=None, total=None):
    # 상단 좌측 네이비 사이드바
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    # 상단 제목 영역
    add_text(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.72), Inches(0.95), Inches(11), Inches(0.35),
                 subtitle, size=14, color=SUB_GRAY)
    # 액센트 라인
    add_rect(slide, Inches(0.7), Inches(1.4), Inches(0.8), Inches(0.06), ACCENT)
    # 푸터
    add_text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
             "Sensor Monitor → 베데스다 iCReaT DCT 서버 전환", size=10, color=SUB_GRAY)
    if page_no is not None and total is not None:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_no} / {total}", size=10, color=SUB_GRAY, align=PP_ALIGN.RIGHT)


def bullet_box(slide, x, y, w, h, items, *, size=16, line_gap=8,
               bullet_color=ACCENT, text_color=GRAY_TEXT, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_gap)
        run = p.add_run()
        run.text = "■  "
        run.font.name = FONT_KR
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = bullet_color
        run2 = p.add_run()
        run2.text = item
        run2.font.name = FONT_KR
        run2.font.size = Pt(size)
        run2.font.color.rgb = text_color
        run2.font.bold = bold_first and i == 0
    return tb


TOTAL = 13


# ============ Slide 1: Title ============
s = prs.slides.add_slide(BLANK)
# 배경
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
# 액센트 가로 라인
add_rect(s, Inches(1.2), Inches(3.55), Inches(1.2), Inches(0.08), ACCENT)
# 부제
add_text(s, Inches(1.2), Inches(2.4), Inches(11), Inches(0.5),
         "2026-05-13 과제 미팅 보고", size=20, color=RGBColor(0xC9, 0xD4, 0xE2))
# 메인 타이틀
add_text(s, Inches(1.2), Inches(2.9), Inches(11), Inches(1.4),
         "Sensor Monitor 앱\n베데스다 iCReaT DCT 서버 전환 진행 현황",
         size=40, bold=True, color=WHITE)
# 보조 설명
add_text(s, Inches(1.2), Inches(4.9), Inches(11), Inches(1.5),
         "기존 CSV 업로드 구조 점검 · 베데스다 회신 정리 · 1차 연동 계획 공유",
         size=18, color=RGBColor(0xB8, 0xC6, 0xD9))
# 하단 정보
add_text(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4),
         "발표자: 이정훈     |     2026-05-13", size=14, color=RGBColor(0x9A, 0xAA, 0xBE))


# ============ Slide 2: 발표 목표 & 핵심 메시지 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "발표 목표", "오늘 공유할 내용의 핵심 메시지", 2, TOTAL)

add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.6),
         "Sensor Monitor 앱의 CSV 업로드 대상을 베데스다 iCReaT DCT 서버로 전환하기 위한",
         size=16, color=GRAY_TEXT)
add_text(s, Inches(0.7), Inches(2.15), Inches(12), Inches(0.6),
         "사전 정리와 협의 진행 상황을 공유합니다.",
         size=16, color=GRAY_TEXT)

# 3개 카드
card_w = Inches(3.95)
card_h = Inches(3.0)
card_y = Inches(3.2)
gap = Inches(0.15)
start_x = Inches(0.7)

messages = [
    ("01", "기존 구조 확인", "기존 앱·서버의 CSV 업로드 구조를\n코드 기준으로 확인했습니다."),
    ("02", "변경사항 정리", "베데스다 측 답변을 바탕으로\n필요한 변경사항을 정리했습니다."),
    ("03", "작업 분리", "바로 구현 가능한 부분과\n베데스다 확인이 필요한 부분을 분리했습니다."),
]
for i, (num, title, body) in enumerate(messages):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, card_y, card_w, card_h, LIGHT_BG)
    add_rect(s, x, card_y, card_w, Inches(0.08), ACCENT)
    add_text(s, x + Inches(0.3), card_y + Inches(0.25), card_w - Inches(0.6), Inches(0.7),
             num, size=36, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), card_y + Inches(1.05), card_w - Inches(0.6), Inches(0.6),
             title, size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), card_y + Inches(1.75), card_w - Inches(0.6), Inches(1.2),
             body, size=14, color=GRAY_TEXT)


# ============ Slide 3: 배경 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "배경", "기존 전송 구조와 이번 작업의 범위", 3, TOTAL)

# 왼쪽 (기존)
add_rect(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(4.6), LIGHT_BG)
add_text(s, Inches(0.95), Inches(2.1), Inches(5.3), Inches(0.5),
         "기존 구조", size=20, bold=True, color=NAVY)
add_rect(s, Inches(0.95), Inches(2.6), Inches(0.6), Inches(0.05), ACCENT)
bullet_box(s, Inches(0.95), Inches(2.85), Inches(5.3), Inches(3.5),
           ["전송 대상: 강원대 서버",
            "전송 방식: multipart/form-data",
            "전송 필드: csvfile, userID,",
            "                  battery, timestamp",
            "서버 구현: Python / Django REST"],
           size=15, line_gap=10)

# 화살표
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                         Inches(6.5), Inches(3.9), Inches(0.5), Inches(0.6))
arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
arr.line.fill.background()

# 오른쪽 (전환)
add_rect(s, Inches(7.1), Inches(1.9), Inches(5.5), Inches(4.6), NAVY)
add_text(s, Inches(7.35), Inches(2.1), Inches(5.1), Inches(0.5),
         "전환 목표", size=20, bold=True, color=WHITE)
add_rect(s, Inches(7.35), Inches(2.6), Inches(0.6), Inches(0.05), ACCENT)
bullet_box(s, Inches(7.35), Inches(2.85), Inches(5.1), Inches(3.5),
           ["전송 대상: 베데스다 iCReaT DCT",
            "전송 방식: multipart/form-data 유지",
            "전송 필드: 5개 (Study/Subject 추가)",
            "인증: 1차 무인증 → 추후 협의",
            "QR 매칭: 업로드 검증 이후 단계"],
           size=15, line_gap=10,
           bullet_color=ACCENT, text_color=WHITE)


# ============ Slide 4: 베데스다 회신 요약 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "베데스다 회신 요약", "전송 필드 정리", 4, TOTAL)

add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5),
         "베데스다 측은 CSV 업로드 엔드포인트를 추가해주기로 답변했습니다.",
         size=16, color=GRAY_TEXT)
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
         "iCReaT는 다중 과제 환경이라 userID 만으로는 부족 → studyId / subjectId 추가 필요.",
         size=16, color=GRAY_TEXT)

# 전송 필드 5개 카드
fields = [
    ("csvfile", "센서 병합 CSV 파일", ACCENT),
    ("battery", "워치 배터리 값", BLUE),
    ("timestamp", "전송 시각", BLUE),
    ("studyId", "과제 식별자 (신규)", NAVY),
    ("subjectId", "피험자 식별자 (신규)", NAVY),
]
card_w = Inches(2.4); card_h = Inches(2.4)
start_x = Inches(0.7); start_y = Inches(3.2)
for i, (name, desc, color) in enumerate(fields):
    x = start_x + (card_w + Inches(0.08)) * i
    add_rect(s, x, start_y, card_w, card_h, LIGHT_BG)
    add_rect(s, x, start_y, card_w, Inches(0.45), color)
    add_text(s, x, start_y + Inches(0.06), card_w, Inches(0.4),
             f"필드 {i+1}", size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, x + Inches(0.15), start_y + Inches(0.7), card_w - Inches(0.3), Inches(0.6),
             name, size=18, bold=True, color=color, align=PP_ALIGN.CENTER,
             font=FONT_MONO)
    add_text(s, x + Inches(0.15), start_y + Inches(1.4), card_w - Inches(0.3), Inches(0.9),
             desc, size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# 하단 메모
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
         "※ 필드명 철자는 베데스다에 보낼 메일에서 다시 확인할 예정",
         size=13, color=SUB_GRAY)


# ============ Slide 5: 인증 방식 정리 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "인증 방식 정리", "1차 목표는 인증 구현이 아닌 업로드 연동", 5, TOTAL)

# 확인했던 방식 3가지
add_text(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.5),
         "사전 확인: 베데스다 서버에 기존 인증 방식이 있을 가능성을 점검했습니다.",
         size=16, color=GRAY_TEXT)

methods = [
    ("Dct-Session-Id 헤더", "세션 기반 인증 가능성", False),
    ("API Key 헤더", "정적 키 인증 가능성", False),
    ("Bearer Token", "토큰 기반 인증 가능성", False),
]
y = Inches(2.7)
for name, desc, ok in methods:
    add_rect(s, Inches(0.7), y, Inches(5.8), Inches(0.6), LIGHT_BG)
    add_text(s, Inches(0.95), y + Inches(0.12), Inches(0.4), Inches(0.4),
             "✕", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.4), y + Inches(0.08), Inches(2.6), Inches(0.5),
             name, size=14, bold=True, color=NAVY)
    add_text(s, Inches(4.0), y + Inches(0.12), Inches(2.4), Inches(0.5),
             desc, size=12, color=SUB_GRAY)
    y += Inches(0.75)

# 결론 박스 (오른쪽)
add_rect(s, Inches(7.0), Inches(2.7), Inches(5.6), Inches(3.6), NAVY)
add_text(s, Inches(7.3), Inches(2.95), Inches(5.0), Inches(0.5),
         "결론", size=20, bold=True, color=WHITE)
add_rect(s, Inches(7.3), Inches(3.45), Inches(0.6), Inches(0.05), ACCENT)
add_text(s, Inches(7.3), Inches(3.7), Inches(5.0), Inches(0.5),
         "현재 베데스다 측에 해당 인증은 구현되어 있지 않음", size=13, color=WHITE)
add_text(s, Inches(7.3), Inches(4.2), Inches(5.0), Inches(0.5),
         "기존 강원대 서버도 무인증 multipart 방식 사용", size=13, color=WHITE)
add_text(s, Inches(7.3), Inches(4.95), Inches(5.0), Inches(1.3),
         "1차 목표는 PPG 포함 센서 데이터\n업로드 연동을 먼저 성공시키는 것",
         size=16, bold=True, color=ACCENT)


# ============ Slide 6: 코드 기준 확인 (Before / After) ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "코드 기준 확인", "OkHttp MultipartBody.FORM 전송 구조", 6, TOTAL)

add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5),
         "모바일 앱은 OkHttp MultipartBody.FORM 으로 CSV 파일과 메타데이터를 함께 전송합니다.",
         size=14, color=GRAY_TEXT)

# Before
add_text(s, Inches(0.7), Inches(2.4), Inches(5.7), Inches(0.4),
         "Before  (강원대 서버 전송)", size=15, bold=True, color=SUB_GRAY)
add_rect(s, Inches(0.7), Inches(2.85), Inches(5.7), Inches(3.4), CODE_BG)
before_code = (
    '.addFormDataPart("csvfile",\n'
    '                  file.name, ...)\n'
    '.addFormDataPart("userID", userID)\n'
    '.addFormDataPart("battery", battery)\n'
    '.addFormDataPart("timestamp",\n'
    '                  timestamp)'
)
add_text(s, Inches(0.9), Inches(3.0), Inches(5.3), Inches(3.2),
         before_code, font=FONT_MONO, size=14, color=CODE_FG)

# Arrow
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                         Inches(6.5), Inches(4.3), Inches(0.4), Inches(0.5))
arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
arr.line.fill.background()

# After
add_text(s, Inches(7.0), Inches(2.4), Inches(5.7), Inches(0.4),
         "After  (베데스다 iCReaT DCT 전송)", size=15, bold=True, color=NAVY)
add_rect(s, Inches(7.0), Inches(2.85), Inches(5.7), Inches(3.4), CODE_BG)
after_code = (
    '.addFormDataPart("csvfile",\n'
    '                  file.name, ...)\n'
    '.addFormDataPart("studyId", studyId)\n'
    '.addFormDataPart("subjectId", subjectId)\n'
    '.addFormDataPart("battery", battery)\n'
    '.addFormDataPart("timestamp",\n'
    '                  timestamp)'
)
# 하이라이트용으로 다른 색
tb = s.shapes.add_textbox(Inches(7.2), Inches(3.0), Inches(5.3), Inches(3.2))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
for i, line in enumerate(after_code.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run(); run.text = line
    run.font.name = FONT_MONO; run.font.size = Pt(14)
    is_new = ("studyId" in line) or ("subjectId" in line)
    run.font.color.rgb = ACCENT if is_new else CODE_FG
    run.font.bold = is_new

# 하단 노트
add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
         "기존 서버 수신부는 Django REST MultiPartParser (request.FILES['csvfile'], request.POST) 구조였습니다.",
         size=12, color=SUB_GRAY)


# ============ Slide 7: 센서 CSV 생성/전송 흐름 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "센서 CSV 생성·전송 흐름", "성공 / 실패 후 파일 이동 로직", 7, TOTAL)

# 흐름 4박스
steps = [
    ("워치 데이터 수집", "센서별 약 5분 단위\nCSV 조각 저장"),
    ("30분 단위 병합", "매시 00분 / 30분에\n조각 CSV를 병합"),
    ("서버 전송", "병합 CSV를\n베데스다 서버로 업로드"),
    ("응답 처리", "HTTP 200 기준\n성공 / 실패 분기"),
]
sw = Inches(2.7); sh = Inches(1.5); sy = Inches(2.0)
for i, (t, b) in enumerate(steps):
    x = Inches(0.7) + (sw + Inches(0.25)) * i
    add_rect(s, x, sy, sw, sh, LIGHT_BG)
    add_rect(s, x, sy, sw, Inches(0.06), BLUE)
    add_text(s, x + Inches(0.15), sy + Inches(0.15), sw - Inches(0.3), Inches(0.4),
             f"STEP {i+1}", size=10, bold=True, color=BLUE)
    add_text(s, x + Inches(0.15), sy + Inches(0.5), sw - Inches(0.3), Inches(0.4),
             t, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.15), sy + Inches(0.95), sw - Inches(0.3), Inches(0.5),
             b, size=11, color=GRAY_TEXT)
    if i < 3:
        ax = x + sw + Inches(0.02)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, sy + Inches(0.6),
                                 Inches(0.2), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
        arr.line.fill.background()

# 성공 / 실패 박스
y = Inches(3.9)
add_rect(s, Inches(0.7), y, Inches(5.9), Inches(2.6), LIGHT_BG)
add_rect(s, Inches(0.7), y, Inches(0.15), Inches(2.6), RGBColor(0x2E, 0x8B, 0x57))
add_text(s, Inches(1.0), y + Inches(0.2), Inches(5.5), Inches(0.5),
         "성공 (HTTP 200)", size=18, bold=True, color=RGBColor(0x2E, 0x8B, 0x57))
bullet_box(s, Inches(1.0), y + Inches(0.85), Inches(5.5), Inches(1.7),
           ["병합 CSV → sended 폴더로 이동",
            "병합에 사용된 원본 조각 CSV 삭제",
            "sended = '서버 전송 성공' 보관 위치"],
           size=13, line_gap=6, bullet_color=RGBColor(0x2E, 0x8B, 0x57))

add_rect(s, Inches(6.8), y, Inches(5.9), Inches(2.6), LIGHT_BG)
add_rect(s, Inches(6.8), y, Inches(0.15), Inches(2.6), ACCENT)
add_text(s, Inches(7.1), y + Inches(0.2), Inches(5.5), Inches(0.5),
         "실패", size=18, bold=True, color=ACCENT)
bullet_box(s, Inches(7.1), y + Inches(0.85), Inches(5.5), Inches(1.7),
           ["병합 CSV 삭제",
            "원본 조각 CSV는 보존",
            "다음 00분 / 30분에 다시 병합·재전송"],
           size=13, line_gap=6, bullet_color=ACCENT)


# ============ Slide 8: Study / Subject ID 처리 계획 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "Study ID / Subject ID 처리 계획", "QR 스캔은 1차 연동 이후 단계", 8, TOTAL)

# 좌측: pbcr_source 분석
add_text(s, Inches(0.7), Inches(1.9), Inches(6), Inches(0.5),
         "pbcr_source 분석 결과", size=18, bold=True, color=NAVY)
add_rect(s, Inches(0.7), Inches(2.5), Inches(0.5), Inches(0.05), ACCENT)
add_text(s, Inches(0.7), Inches(2.8), Inches(6), Inches(0.5),
         "기존 앱은 QR 코드 JSON에서 식별자를 읽어 매칭하는 구조였습니다.",
         size=13, color=GRAY_TEXT)

add_rect(s, Inches(0.7), Inches(3.4), Inches(5.8), Inches(2.3), CODE_BG)
add_text(s, Inches(0.95), Inches(3.55), Inches(5.4), Inches(0.5),
         "QR JSON 매핑", size=12, bold=True, color=ACCENT, font=FONT_MONO)
add_text(s, Inches(0.95), Inches(3.95), Inches(5.4), Inches(0.55),
         'stdy_no       →   Study ID',
         size=15, font=FONT_MONO, color=CODE_FG)
add_text(s, Inches(0.95), Inches(4.55), Inches(5.4), Inches(0.55),
         'subject_id    →   Subject ID',
         size=15, font=FONT_MONO, color=CODE_FG)
add_text(s, Inches(0.95), Inches(5.2), Inches(5.4), Inches(0.5),
         "Sensor Monitor 에는 QR 스캔 기능이 아직 없음",
         size=11, color=RGBColor(0xC9, 0xD4, 0xE2), font=FONT_MONO)

# 우측: 1차 진행 계획
add_text(s, Inches(7.0), Inches(1.9), Inches(6), Inches(0.5),
         "1차 진행 계획", size=18, bold=True, color=NAVY)
add_rect(s, Inches(7.0), Inches(2.5), Inches(0.5), Inches(0.05), ACCENT)

plans = [
    ("STEP 1", "베데스다에 테스트용\nStudy ID / Subject ID 1쌍 요청"),
    ("STEP 2", "앱에 임시 하드코딩 후\n업로드 검증 진행"),
    ("STEP 3", "업로드 연동 안정화 후\nQR 스캔 기능 추가"),
]
py = Inches(2.85)
for tag, body in plans:
    add_rect(s, Inches(7.0), py, Inches(5.6), Inches(1.0), LIGHT_BG)
    add_rect(s, Inches(7.0), py, Inches(1.1), Inches(1.0), NAVY)
    add_text(s, Inches(7.0), py + Inches(0.3), Inches(1.1), Inches(0.5),
             tag, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.2), py + Inches(0.18), Inches(4.3), Inches(0.85),
             body, size=13, color=GRAY_TEXT)
    py += Inches(1.1)

add_text(s, Inches(7.0), Inches(6.25), Inches(5.6), Inches(0.5),
         "※ QR 스캔: 카메라 권한 / 라이브러리 / UI 수정 필요 → 후순위",
         size=11, color=SUB_GRAY)


# ============ Slide 9: 샘플 데이터 준비 현황 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "샘플 데이터 준비 현황", "Sensor_monitor/sensor_data.zip", 9, TOTAL)

add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.5),
         "베데스다 엔드포인트 테스트용 샘플 데이터를 zip 으로 정리했습니다.",
         size=15, color=GRAY_TEXT)

# 좌측 통계
stats = [
    ("175",  "전체 엔트리 수"),
    ("172",  "CSV 파일 수"),
    ("143",  "sended 30분 병합 CSV"),
    ("3",    "로그 파일 수"),
]
sx = Inches(0.7); sy = Inches(2.7)
sw = Inches(2.85); sh = Inches(1.7)
for i, (n, label) in enumerate(stats):
    col = i % 2; row = i // 2
    x = sx + (sw + Inches(0.15)) * col
    y = sy + (sh + Inches(0.2)) * row
    add_rect(s, x, y, sw, sh, LIGHT_BG)
    add_text(s, x, y + Inches(0.3), sw, Inches(0.8),
             n, size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.15), sw, Inches(0.4),
             label, size=12, color=SUB_GRAY, align=PP_ALIGN.CENTER)

# 우측 센서 목록
add_rect(s, Inches(6.85), Inches(2.7), Inches(5.8), Inches(3.85), NAVY)
add_text(s, Inches(7.15), Inches(2.9), Inches(5.4), Inches(0.5),
         "포함 센서 (7종)", size=18, bold=True, color=WHITE)
add_rect(s, Inches(7.15), Inches(3.4), Inches(0.5), Inches(0.05), ACCENT)

sensors = ["PpgGreen", "HeartRate", "Accelerometer",
           "Gyroscope", "Gravity", "Light", "StepCount"]
sy2 = Inches(3.6)
for i, name in enumerate(sensors):
    col = i % 2; row = i // 2
    x = Inches(7.15) + Inches(2.7) * col
    y = sy2 + Inches(0.42) * row
    add_text(s, x, y, Inches(0.3), Inches(0.4),
             "▸", size=14, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.3), y, Inches(2.3), Inches(0.4),
             name, size=14, color=WHITE, font=FONT_MONO)

add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
         "각 센서 폴더의 sended 하위 CSV 를 베데스다 엔드포인트 테스트 샘플로 활용 예정",
         size=12, color=SUB_GRAY)


# ============ Slide 10: 현재 정리한 산출물 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "현재 정리한 산출물", "협의·분석을 위해 정리한 문서", 10, TOTAL)

docs = [
    ("01", "bethesda_email_draft.md",
     "베데스다 담당자에게 보낼 회신 메일 초안"),
    ("02", "bethesda_multipart_reference.md",
     "multipart 전송 / 수신 및 샘플 데이터 부연 설명"),
    ("03", "bethesda_csv_upload_reply_analysis.md",
     "베데스다 회신 분석 및 수정 포인트 보고서"),
]
y = Inches(2.0)
for num, name, desc in docs:
    add_rect(s, Inches(0.7), y, Inches(12), Inches(1.1), LIGHT_BG)
    add_rect(s, Inches(0.7), y, Inches(0.15), Inches(1.1), ACCENT)
    add_text(s, Inches(1.0), y + Inches(0.18), Inches(0.8), Inches(0.7),
             num, size=24, bold=True, color=ACCENT)
    add_text(s, Inches(1.9), y + Inches(0.18), Inches(10), Inches(0.5),
             name, size=18, bold=True, color=NAVY, font=FONT_MONO)
    add_text(s, Inches(1.9), y + Inches(0.65), Inches(10), Inches(0.4),
             desc, size=13, color=GRAY_TEXT)
    y += Inches(1.25)

# 코드 정리
add_rect(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.85), NAVY)
add_text(s, Inches(1.0), Inches(6.0), Inches(11.5), Inches(0.5),
         "코드", size=14, bold=True, color=ACCENT)
add_text(s, Inches(1.0), Inches(6.3), Inches(11.5), Inches(0.4),
         "CSV 전송 성공 / 실패 후 파일 이동 로직 정리 완료",
         size=14, color=WHITE)


# ============ Slide 11: 다음 진행 계획 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "다음 진행 계획", "8단계로 정리한 다음 액션", 11, TOTAL)

plans = [
    "베데스다에 회신 메일 발송",
    "신규 업로드 엔드포인트 URL · 응답 형식 수신",
    "테스트용 Study ID / Subject ID 수신",
    "Sensor Monitor 앱의 전송 URL 과 필드 수정",
    "curl / Postman 으로 단독 업로드 테스트",
    "앱 디버그 빌드로 00분 / 30분 주기 업로드 테스트",
    "베데스다 담당자와 서버 적재 여부 확인",
    "QR 스캔 또는 수동 입력 방식으로 Study/Subject 매칭 추가",
]
# 2열로 배치
col_w = Inches(5.95); row_h = Inches(1.0)
start_x = Inches(0.7); start_y = Inches(1.95)
for i, p in enumerate(plans):
    col = i % 2; row = i // 2
    x = start_x + (col_w + Inches(0.15)) * col
    y = start_y + (row_h + Inches(0.12)) * row
    add_rect(s, x, y, col_w, row_h, LIGHT_BG)
    # 번호 원
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                x + Inches(0.2), y + Inches(0.2),
                                Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    tf = circle.text_frame; tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = str(i + 1)
    run.font.name = FONT_KR; run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = WHITE
    add_text(s, x + Inches(0.95), y + Inches(0.28), col_w - Inches(1.1), Inches(0.5),
             p, size=13, color=GRAY_TEXT)


# ============ Slide 12: 미팅 확인 요청 ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "미팅에서 확인받고 싶은 사항", "방향 합의가 필요한 5가지 질문", 12, TOTAL)

qs = [
    "1차 연동은 무인증 방식으로 진행해도 되는지",
    "Study ID / Subject ID 를 임시 하드코딩해서 먼저 테스트해도 되는지",
    "QR 스캔 기능은 1차 업로드 연동 이후 후순위로 두어도 되는지",
    "베데스다 측에 샘플 CSV 전체 ZIP 을 공유해도 되는지",
    "서버 성공 응답은 HTTP 200 기준으로 요청해도 되는지",
]
y = Inches(2.0)
for i, q in enumerate(qs):
    add_rect(s, Inches(0.7), y, Inches(12), Inches(0.85), LIGHT_BG)
    # Q 라벨
    add_rect(s, Inches(0.7), y, Inches(0.8), Inches(0.85), NAVY)
    add_text(s, Inches(0.7), y + Inches(0.2), Inches(0.8), Inches(0.5),
             f"Q{i+1}", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y + Inches(0.22), Inches(10.8), Inches(0.5),
             q, size=14, color=GRAY_TEXT)
    y += Inches(1.0)


# ============ Slide 13: 예상 Q&A ============
s = prs.slides.add_slide(BLANK)
header_bar(s, "예상 질문 정리", "발표 중·후 응답 준비", 13, TOTAL)

qas = [
    ("왜 QR 스캔을 바로 구현하지 않나요?",
     "카메라 권한·라이브러리·UI 수정이 필요. 1차는 엔드포인트 연동 검증이\n우선이라 하드코딩된 테스트 ID 로 진행 후 QR 추가가 리스크가 낮습니다."),
    ("sended 폴더의 의미는?",
     "서버 전송에 성공한 30분 단위 병합 CSV 보관 위치.\n실패 시에는 병합본 삭제 + 원본 조각 보존으로 재전송 가능하게 정리."),
    ("성공 응답 기준은?",
     "앱 코드는 OkHttp response.isSuccessful (HTTP 200~299) 기준.\n베데스다에는 성공 시 HTTP 200 반환을 요청할 예정."),
    ("battery 값은 문자열 / 숫자?",
     "워치에서 받은 숫자값을 multipart form field 로 전송. 서버에서는\n숫자로 파싱하여 저장하면 됨."),
]
col_w = Inches(5.95); row_h = Inches(2.2)
start_x = Inches(0.7); start_y = Inches(1.95)
for i, (q, a) in enumerate(qas):
    col = i % 2; row = i // 2
    x = start_x + (col_w + Inches(0.15)) * col
    y = start_y + (row_h + Inches(0.15)) * row
    add_rect(s, x, y, col_w, row_h, LIGHT_BG)
    add_rect(s, x, y, Inches(0.15), row_h, ACCENT)
    add_text(s, x + Inches(0.35), y + Inches(0.15), col_w - Inches(0.5), Inches(0.45),
             "Q", size=14, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.65), y + Inches(0.13), col_w - Inches(0.8), Inches(0.5),
             q, size=13, bold=True, color=NAVY)
    add_text(s, x + Inches(0.35), y + Inches(0.85), col_w - Inches(0.5), Inches(0.45),
             "A", size=14, bold=True, color=BLUE)
    add_text(s, x + Inches(0.65), y + Inches(0.85), col_w - Inches(0.8), Inches(1.2),
             a, size=12, color=GRAY_TEXT)


import os
# 산출물은 reports/decks/에 저장(builders/ 기준 상대경로). 저장소 이동에도 안전. (경로 버그 수정 2026-06-29, 폴더정리 2026-06-29)
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "decks", "2026-05-13_meeting_update.pptx")
prs.save(out)
print(f"saved: {out}")
