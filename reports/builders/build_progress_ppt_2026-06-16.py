# -*- coding: utf-8 -*-
"""2026-06-17(수) 발표용 진행상황 보고 PPT 생성.
목적: 직전 보고(2026-05-27) 이후의 경과 — 베데스다 API 정식 수신 → 엔드포인트 단독 검증(DB 적재) 완료.
말투: ~임 / ~했음 (개조식).
스타일: 05-13 / 05-27 덱과 동일 팔레트/폰트.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import os
# 산출물은 reports/decks/에 저장(builders/ 기준 상대경로). 저장소 이동에도 안전. (경로 버그 수정 2026-06-29, 폴더정리 2026-06-29)
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "decks", "[ICreaT과제] 진행상황 보고서_2026-06-17.pptx")

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xE8, 0x74, 0x22)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
SUB_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE0, 0xE6, 0xF0)
DIM = RGBColor(0x8A, 0x95, 0xA8)

FONT_KR = "맑은 고딕"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def add_rect(s, x, y, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(s, x, y, w, h, lines, *, font=FONT_KR, size=18, bold=False,
             color=GRAY_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=2):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(line, tuple):
            text, ov = line
        else:
            text, ov = line, {}
        run = p.add_run()
        run.text = text
        run.font.name = ov.get("font", font)
        run.font.size = Pt(ov.get("size", size))
        run.font.bold = ov.get("bold", bold)
        run.font.color.rgb = ov.get("color", color)
    return tb


def header(s, idx, total, title, subtitle):
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_text(s, Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.72), Inches(0.98), Inches(11.8), Inches(0.35),
             subtitle, size=14, color=SUB_GRAY)
    add_rect(s, Inches(0.7), Inches(1.42), Inches(0.8), Inches(0.06), ACCENT)
    add_text(s, Inches(0.7), Inches(7.08), Inches(8), Inches(0.3),
             "Sensor Monitor → 베데스다 iCReaT DCT 서버 전환", size=10, color=SUB_GRAY)
    add_text(s, Inches(11.4), Inches(7.08), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=10, color=SUB_GRAY, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, *, gap=0.52, size=14, dot_color=ACCENT,
            head_color=NAVY, body_color=GRAY_TEXT):
    cy = y
    for head, body in items:
        add_text(s, x, cy, Inches(0.25), Inches(0.3), "▸", size=size, bold=True, color=dot_color)
        add_text(s, x + Inches(0.32), cy - Inches(0.02), w - Inches(0.32), Inches(0.3),
                 head, size=size, bold=True, color=head_color)
        if body:
            add_text(s, x + Inches(0.32), cy + Inches(0.26), w - Inches(0.32), Inches(0.3),
                     body, size=size - 3, color=body_color)
        cy += Inches(gap)


TOTAL = 10

# ============ 1. 타이틀 ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(5.0), SLIDE_W, Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0),
         "ICReaT 과제 진행상황 보고", size=40, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.6),
         "Sensor Monitor → 베데스다 iCReaT DCT 서버 CSV 업로드 전환", size=18, color=RGBColor(0xC8, 0xD4, 0xE4))
add_text(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
         "2026-06-17 (수)  ·  직전 보고(2026-05-27) 이후 경과 중심", size=14, color=RGBColor(0x9F, 0xB0, 0xC8))

# ============ 2. 직전(05-27) 시점 요약 ============
s = slide()
header(s, 2, TOTAL, "직전(05-27) 시점", "코드 1차 수정은 끝났고, 두 값만 받으면 테스트 가능한 상태였음")
add_rect(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.06), DIM)
bullets(s, Inches(0.9), Inches(2.1), Inches(11.4), [
    ("앱 전송 규격을 베데스다 형식으로 1차 전환 완료", "userID → studyId/subjectId, 5개 파일 수정"),
    ("로그인 단계는 무인증 정책에 맞춰 우회 처리", "네트워크 login() 제거 방향"),
    ("단, 실제 검증은 두 값에 막혀 있었음", "엔드포인트 URL · 테스트용 Study/Subject ID 미수신"),
    ("코드에는 placeholder만 들어가 있었음", "REQUEST_URL = .../REPLACE_WITH_BETHESDA_ENDPOINT"),
], gap=0.78, size=16)
add_rect(s, Inches(0.9), Inches(6.2), Inches(11.4), Inches(0.7), LIGHT_BG)
add_rect(s, Inches(0.9), Inches(6.2), Inches(0.1), Inches(0.7), ACCENT)
add_text(s, Inches(1.2), Inches(6.37), Inches(11), Inches(0.4),
         "→ 즉, 05-27엔 '코드 준비 완료, 서버 값 대기' 상태 — 실제 전송은 한 번도 안 해본 시점이었음",
         size=14, bold=True, color=NAVY)

# ============ 3. 그 이후 경과 (타임라인) ============
s = slide()
header(s, 3, TOTAL, "05-27 이후 경과", "API 가이드 수신 → 실값 반영 → 경로 이슈 해결 → 단독 검증 → 적재 확인 요청")
steps = [
    ("1", "베데스다 API 연동 가이드 정식 수신", "엔드포인트·5필드·응답코드·테스트 대상자·DB 저장규격 문서화됨", BLUE),
    ("2", "실제 값 코드 반영", "placeholder → 실 엔드포인트, 테스트 ID(C250005 / 121-001) 적용", BLUE),
    ("3", "경로 이슈 발견·해결", "가이드 경로가 404(200+HTML) → 동작 경로(루트 /invoke/...)로 교정", ACCENT),
    ("4", "curl 단독 전송 검증 통과", "HTTP 200 + status:success 응답 확인", GREEN),
    ("5", "서버 DB 적재 확인 요청 (진행 중)", "응답은 성공이나 실제 DB 적재는 직접 확인 불가 → 담당자에 확인 요청함", ACCENT),
]
ty = Inches(2.0)
for num, head, body, color in steps:
    add_rect(s, Inches(0.9), ty, Inches(0.55), Inches(0.55), color)
    add_text(s, Inches(0.9), ty + Inches(0.06), Inches(0.55), Inches(0.45),
             num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.65), ty - Inches(0.02), Inches(10.8), Inches(0.4),
             head, size=17, bold=True, color=NAVY)
    add_text(s, Inches(1.65), ty + Inches(0.34), Inches(10.8), Inches(0.35),
             body, size=13, color=SUB_GRAY)
    ty += Inches(0.96)

# ============ 4. 핵심 성과 (강조) ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(1.0), Inches(0.9), Inches(0.8), Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(1.1), Inches(11.3), Inches(0.6),
         "이번 보고 핵심 성과", size=30, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.9),
         "베데스다 서버 엔드포인트 전송 성공 응답 확인",
         size=26, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(2.95), Inches(11.3), Inches(0.5),
         "센서 CSV 1건 전송 → 200 + status:success 응답 확인 (DB 적재 여부는 담당자 확인 대기)", size=15, color=RGBColor(0xC8,0xD4,0xE4))
cards = [
    ("200", "HTTP 응답 성공", GREEN),
    ("success", "JSON status 정상", GREEN),
    ("적재 확인 중", "담당자 회신 대기", ACCENT),
]
cx = Inches(1.0)
for big, desc, color in cards:
    add_rect(s, cx, Inches(3.9), Inches(3.7), Inches(1.9), RGBColor(0x16, 0x33, 0x55))
    add_rect(s, cx, Inches(3.9), Inches(3.7), Inches(0.12), color)
    add_text(s, cx, Inches(4.35), Inches(3.7), Inches(0.8),
             big, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)
    add_text(s, cx, Inches(5.25), Inches(3.7), Inches(0.4),
             desc, size=14, color=RGBColor(0xC8,0xD4,0xE4), align=PP_ALIGN.CENTER)
    cx += Inches(4.05)
add_text(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.5),
         "→ 05-27의 '대기' 상태를 넘어, 전송이 성공 응답을 받는 단계까지 도달함 (실제 DB 적재는 담당자 확인 후 확정)",
         size=14, bold=True, color=ACCENT)

# ============ 5. 받은 API 명세 요약 ============
s = slide()
header(s, 5, TOTAL, "수신한 API 명세", "베데스다 제공 연동 가이드 — 검증의 전제")
add_rect(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(0.55), CODE_BG)
add_text(s, Inches(1.1), Inches(1.96), Inches(11.2), Inches(0.4),
         "POST https://icreatdct.btsd.io/invoke/DCT/Sensor/uploadCsv   (multipart, 무인증, 성공=200)",
         size=14, bold=True, color=CODE_FG, font=FONT_MONO)
# 좌: 응답 코드 표
add_text(s, Inches(0.9), Inches(2.65), Inches(5.6), Inches(0.4),
         "응답 코드", size=15, bold=True, color=NAVY)
codes = [
    ("200", "성공 — 업로드 완료", GREEN),
    ("401", "동일 파일명 중복 전송", ACCENT),
    ("403", "필수 파라미터 누락/오류", ACCENT),
    ("405", "등록되지 않은 대상자", RED),
    ("406", "DB 접속/처리 오류", RED),
]
cy = Inches(3.15)
for code, desc, color in codes:
    add_rect(s, Inches(0.9), cy, Inches(1.0), Inches(0.5), color)
    add_text(s, Inches(0.9), cy + Inches(0.07), Inches(1.0), Inches(0.4),
             code, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)
    add_text(s, Inches(2.05), cy + Inches(0.1), Inches(4.4), Inches(0.4),
             desc, size=13, color=GRAY_TEXT)
    cy += Inches(0.62)
# 우: 테스트 대상자 + 저장
add_text(s, Inches(7.0), Inches(2.65), Inches(5.5), Inches(0.4),
         "테스트 대상자 (가이드 §8)", size=15, bold=True, color=NAVY)
subs = [("C250005", "121-001  ← 검증에 사용"), ("C250002", "001-002"), ("C250002", "001-001")]
sy = Inches(3.15)
for st, sub in subs:
    add_rect(s, Inches(7.0), sy, Inches(5.4), Inches(0.5), LIGHT_BG)
    add_text(s, Inches(7.2), sy + Inches(0.1), Inches(2.0), Inches(0.4),
             st, size=14, bold=True, color=NAVY, font=FONT_MONO)
    add_text(s, Inches(9.2), sy + Inches(0.12), Inches(3.1), Inches(0.4),
             sub, size=13, color=GRAY_TEXT, font=FONT_MONO)
    sy += Inches(0.62)
add_text(s, Inches(7.0), Inches(5.15), Inches(5.5), Inches(0.4),
         "서버 저장", size=15, bold=True, color=NAVY)
add_text(s, Inches(7.0), Inches(5.6), Inches(5.5), Inches(1.0), [
    ("원본 CSV: sensor_upload/{study}/{subject}/{날짜}/", {"size": 12, "color": SUB_GRAY, "font": FONT_MONO}),
    ("이력 DB:  dct_sensor_upload", {"size": 12, "color": SUB_GRAY, "font": FONT_MONO}),
], space_after=4)

# ============ 6. 경로 트러블슈팅 ============
s = slide()
header(s, 6, TOTAL, "검증 중 발견·해결한 이슈", "가이드 경로가 그대로는 동작하지 않았음 — 진단 후 교정")
# before (가이드 경로)
add_rect(s, Inches(0.9), Inches(2.0), Inches(5.6), Inches(2.0), CODE_BG)
add_text(s, Inches(1.15), Inches(2.15), Inches(5.1), Inches(0.4),
         "가이드 문서 경로", size=14, bold=True, color=DIM)
add_text(s, Inches(1.15), Inches(2.6), Inches(5.2), Inches(1.3), [
    ("/iCReaT_DCT/invoke", {"color": RGBColor(0xE0,0x6C,0x6C)}),
    ("        /DCT/Sensor/uploadCsv", {"color": RGBColor(0xE0,0x6C,0x6C)}),
    ("", {}),
    ("→ HTTP 200 이지만 HTML 에러", {"color": DIM, "size": 13}),
    ("   페이지 반환 (실질 404)", {"color": DIM, "size": 13}),
], size=14, font=FONT_MONO, space_after=3)
add_text(s, Inches(6.6), Inches(2.7), Inches(0.7), Inches(0.6),
         "→", size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
# after (동작 경로)
add_rect(s, Inches(7.3), Inches(2.0), Inches(5.2), Inches(2.0), CODE_BG)
add_text(s, Inches(7.55), Inches(2.15), Inches(4.7), Inches(0.4),
         "실제 동작 경로", size=14, bold=True, color=GREEN)
add_text(s, Inches(7.55), Inches(2.6), Inches(4.8), Inches(1.3), [
    ("/invoke", {"color": RGBColor(0x8B,0xD4,0x9B)}),
    ("    /DCT/Sensor/uploadCsv", {"color": RGBColor(0x8B,0xD4,0x9B)}),
    ("", {}),
    ("→ 200 + status:success", {"color": RGBColor(0x8B,0xD4,0x9B), "size": 13}),
    ("   DB 적재 확인", {"color": RGBColor(0x8B,0xD4,0x9B), "size": 13}),
], size=14, font=FONT_MONO, space_after=3)
# 대응
add_text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.4),
         "코드 측 대응", size=15, bold=True, color=NAVY)
bullets(s, Inches(0.9), Inches(5.0), Inches(11.4), [
    ("성공 판정을 HTTP 코드만으로 하지 않도록 보강", "잘못된 경로도 200+HTML을 주기 때문"),
    ("응답 body의 JSON status==\"success\"까지 확인하는 가드 추가", "isSuccessBody() 신설"),
    ("실패 시 응답 body를 로그에 남겨 원인 추적 가능하게 함", ""),
], gap=0.55, size=14)

# ============ 7. 05-27 대비 추가 수정한 코드 ============
s = slide()
header(s, 7, TOTAL, "05-27 대비 추가로 바뀐 코드", "placeholder를 실값으로, 그리고 성공 판정 보강")
rows = [
    ("ServerConnection.kt", "REQUEST_URL을 실제 동작 엔드포인트로 확정, isSuccessBody() 추가", GREEN),
    ("ServerConnection.kt", "login() 완전 제거 → enterSensor()로 대체 (무인증 진입)", BLUE),
    ("DeviceInfo.kt", "테스트 Study/Subject를 실제 제공값(C250005 / 121-001)으로 설정", GREEN),
    ("LoginActivity.kt", "네트워크 로그인·자동로그인 호출 제거, 바로 진입", BLUE),
    ("AcceptService / SensorActivity", "studyId/subjectId 2값 기준으로 호출·초기화 정리", BLUE),
]
ry = Inches(2.05)
for fname, desc, color in rows:
    add_rect(s, Inches(0.9), ry, Inches(4.2), Inches(0.82), LIGHT_BG)
    add_rect(s, Inches(0.9), ry, Inches(0.1), Inches(0.82), color)
    add_text(s, Inches(1.15), ry + Inches(0.21), Inches(3.9), Inches(0.45),
             fname, size=13, bold=True, color=NAVY, font=FONT_MONO)
    add_text(s, Inches(5.3), ry + Inches(0.13), Inches(7.2), Inches(0.6),
             desc, size=13, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    ry += Inches(0.92)

# ============ 8. 05-27 → 06-17 상태 비교 ============
s = slide()
header(s, 8, TOTAL, "05-27 → 06-17 상태 비교", "대기 항목 대부분 검증 완료 — DB 적재 확인만 진행 중")
add_rect(s, Inches(0.9), Inches(2.0), Inches(3.9), Inches(0.55), NAVY)
add_rect(s, Inches(4.85), Inches(2.0), Inches(3.75), Inches(0.55), RGBColor(0x88,0x88,0x88))
add_rect(s, Inches(8.65), Inches(2.0), Inches(3.85), Inches(0.55), GREEN)
add_text(s, Inches(0.9), Inches(2.1), Inches(3.9), Inches(0.4), "항목", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.85), Inches(2.1), Inches(3.75), Inches(0.4), "05-27", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.65), Inches(2.1), Inches(3.85), Inches(0.4), "06-17", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
comp = [
    ("엔드포인트 URL", "미수신 (placeholder)", "확정·교정 완료"),
    ("테스트 Study/Subject", "미수신", "C250005 / 121-001 적용"),
    ("응답 규격", "가정(200)", "200/401/403/405/406 확인"),
    ("단독 전송 검증", "미실시", "통과 (200+success)"),
]
comp_pending = [
    ("서버 DB 적재", "미확인", "담당자 확인 요청 중"),
]
ry = Inches(2.62)
for item, before, after in comp:
    add_rect(s, Inches(0.9), ry, Inches(3.9), Inches(0.66), LIGHT_BG)
    add_rect(s, Inches(4.85), ry, Inches(3.75), Inches(0.66), RGBColor(0xF0,0xE6,0xE0))
    add_rect(s, Inches(8.65), ry, Inches(3.85), Inches(0.66), RGBColor(0xE3,0xF1,0xE8))
    add_text(s, Inches(1.1), ry + Inches(0.16), Inches(3.6), Inches(0.4), item, size=13, bold=True, color=NAVY)
    add_text(s, Inches(5.0), ry + Inches(0.16), Inches(3.5), Inches(0.4), before, size=12, color=SUB_GRAY)
    add_text(s, Inches(8.8), ry + Inches(0.16), Inches(3.6), Inches(0.4), [("✔ " + after, {"color": GREEN, "bold": True, "size": 12})], size=12)
    ry += Inches(0.74)
for item, before, after in comp_pending:
    add_rect(s, Inches(0.9), ry, Inches(3.9), Inches(0.66), LIGHT_BG)
    add_rect(s, Inches(4.85), ry, Inches(3.75), Inches(0.66), RGBColor(0xF0,0xE6,0xE0))
    add_rect(s, Inches(8.65), ry, Inches(3.85), Inches(0.66), RGBColor(0xFB,0xF0,0xE2))
    add_text(s, Inches(1.1), ry + Inches(0.16), Inches(3.6), Inches(0.4), item, size=13, bold=True, color=NAVY)
    add_text(s, Inches(5.0), ry + Inches(0.16), Inches(3.5), Inches(0.4), before, size=12, color=SUB_GRAY)
    add_text(s, Inches(8.8), ry + Inches(0.16), Inches(3.6), Inches(0.4), [("… " + after, {"color": ACCENT, "bold": True, "size": 12})], size=12)
    ry += Inches(0.74)

# ============ 9. 앞으로 할 것 ============
s = slide()
header(s, 9, TOTAL, "앞으로 할 것", "앱 종단간 → QR 스캔 → 인증")
phases = [
    ("다음", "앱 종단간 테스트 — 워치→앱→서버 30분 사이클 자동 전송 및 다중 센서/연속 전송 확인함", ACCENT),
    ("Phase 1", "QR 스캔 기능 추가 — 카메라로 Study/Subject 자동 입력, 테스트 하드코딩 제거함", BLUE),
    ("후속", "인증(세션/API Key) 및 중복·오류코드(401/405) 운영 대응은 업로드 안정화 후 논의함", NAVY),
]
py2 = Inches(2.2)
for tag, body, color in phases:
    add_rect(s, Inches(0.9), py2, Inches(2.2), Inches(1.1), color)
    add_text(s, Inches(0.9), py2 + Inches(0.36), Inches(2.2), Inches(0.5),
             tag, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.25), py2, Inches(9.25), Inches(1.1), LIGHT_BG)
    add_text(s, Inches(3.55), py2 + Inches(0.18), Inches(8.8), Inches(0.8),
             body, size=14, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    py2 += Inches(1.4)

# ============ 10. 요약 ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(1.0), Inches(1.0), Inches(0.8), Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.7),
         "요약", size=32, bold=True, color=WHITE)
summary = [
    "05-27엔 코드 준비만 끝났고, 서버 값(엔드포인트·테스트ID) 대기 상태였음",
    "이후 베데스다 API 가이드를 정식 수신하고 실값을 코드에 반영했음",
    "가이드 경로가 그대로는 동작하지 않아 진단·교정했음 (루트 /invoke 경로)",
    "단독 전송 검증 통과 — 200 + status:success 응답 확인함 (DB 적재는 담당자 확인 요청 중)",
    "다음은 적재 확인 회신 받고, 앱 종단간(30분 사이클) 테스트 → QR 스캔·인증 순서로 진행함",
]
sy2 = Inches(2.5)
for line in summary:
    add_text(s, Inches(1.0), sy2, Inches(0.3), Inches(0.4), "▸", size=18, bold=True, color=ACCENT)
    add_text(s, Inches(1.4), sy2, Inches(11.0), Inches(0.6), line, size=18, color=WHITE)
    sy2 += Inches(0.78)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides))
