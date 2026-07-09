# -*- coding: utf-8 -*-
"""2026-07-01(수) 발표용 진행상황 보고 PPT 생성.
목적: 직전 보고(2026-06-17) 이후의 경과 — 종단간 자동 전송 + 베데스다 담당자 수신 확인 완료(적재 미결 종료).
말투: ~임 / ~했음 (개조식).
스타일: 05-13 / 05-27 / 06-17 덱과 동일 팔레트/폰트.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import os
# 산출물은 reports/decks/에 저장(builders/ 기준 상대경로). 저장소 이동에도 안전.
# [2026-07-01] 이유: 07-01 회차 발표자료 신규 | 목적: 06-17 이후 경과(수신 확인) 보고
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "decks", "[ICreaT과제] 진행상황 보고서_2026-07-01.pptx")

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xE8, 0x74, 0x22)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
SUB_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE0, 0xE6, 0xF0)
DIM = RGBColor(0x8A, 0x95, 0xA8)

FONT_KR = "맑은 고딕"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def add_rect(s, x, y, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(s, x, y, w, h, lines, *, font=FONT_KR, size=18, bold=False,
             color=GRAY_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=2):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(line, tuple):
            text, ov = line
        else:
            text, ov = line, {}
        run = p.add_run()
        run.text = text
        run.font.name = ov.get("font", font)
        run.font.size = Pt(ov.get("size", size))
        run.font.bold = ov.get("bold", bold)
        run.font.color.rgb = ov.get("color", color)
    return tb


def header(s, idx, total, title, subtitle):
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_text(s, Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.72), Inches(0.98), Inches(11.8), Inches(0.35),
             subtitle, size=14, color=SUB_GRAY)
    add_rect(s, Inches(0.7), Inches(1.42), Inches(0.8), Inches(0.06), ACCENT)
    add_text(s, Inches(0.7), Inches(7.08), Inches(8), Inches(0.3),
             "Sensor Monitor → 베데스다 iCReaT DCT 서버 전환", size=10, color=SUB_GRAY)
    add_text(s, Inches(11.4), Inches(7.08), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=10, color=SUB_GRAY, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, *, gap=0.52, size=14, dot_color=ACCENT,
            head_color=NAVY, body_color=GRAY_TEXT):
    cy = y
    for head, body in items:
        add_text(s, x, cy, Inches(0.25), Inches(0.3), "▸", size=size, bold=True, color=dot_color)
        add_text(s, x + Inches(0.32), cy - Inches(0.02), w - Inches(0.32), Inches(0.3),
                 head, size=size, bold=True, color=head_color)
        if body:
            add_text(s, x + Inches(0.32), cy + Inches(0.26), w - Inches(0.32), Inches(0.3),
                     body, size=size - 3, color=body_color)
        cy += Inches(gap)


TOTAL = 9

# ============ 1. 타이틀 ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(5.0), SLIDE_W, Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0),
         "ICReaT 과제 진행상황 보고", size=40, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.6),
         "Sensor Monitor → 베데스다 iCReaT DCT 서버 CSV 업로드 전환", size=18, color=RGBColor(0xC8, 0xD4, 0xE4))
add_text(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
         "2026-07-01 (수)  ·  직전 보고(2026-06-17) 이후 경과 중심", size=14, color=RGBColor(0x9F, 0xB0, 0xC8))

# ============ 2. 직전(06-17) 시점 요약 ============
s = slide()
header(s, 2, TOTAL, "직전(06-17) 시점", "단독 1건 전송은 성공했으나, 실제 서버 적재 확인이 미결이었음")
add_rect(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.06), DIM)
bullets(s, Inches(0.9), Inches(2.1), Inches(11.4), [
    ("베데스다 API 가이드 정식 수신·반영 완료", "실 엔드포인트 + 테스트 ID(C250005 / 121-001)"),
    ("가이드 경로 404 이슈 진단·교정", "루트 /invoke/DCT/Sensor/uploadCsv 로 동작 확인"),
    ("curl 단독 전송 검증 통과", "HTTP 200 + status:success 응답 확인"),
    ("단, 실제 서버 DB 적재 여부는 미확인이었음", "응답은 성공이나 적재는 직접 확인 불가 → 담당자 확인 요청 중"),
], gap=0.78, size=16)
add_rect(s, Inches(0.9), Inches(6.2), Inches(11.4), Inches(0.7), LIGHT_BG)
add_rect(s, Inches(0.9), Inches(6.2), Inches(0.1), Inches(0.7), ACCENT)
add_text(s, Inches(1.2), Inches(6.37), Inches(11), Inches(0.4),
         "→ 즉, 06-17엔 '단독 1건 성공, 서버 수신·적재 확인 대기' 상태였음 — 이번 회차에 이 미결을 닫음",
         size=14, bold=True, color=NAVY)

# ============ 3. 그 이후 경과 (타임라인) ============
s = slide()
header(s, 3, TOTAL, "06-17 이후 경과", "종단간 자동 전송 → 베데스다 수신 확인 회신 → 미결 종료")
steps = [
    ("1", "종단간 자동 전송 실시", "워치→앱→서버 30분 사이클로 06-30 15:00 / 15:30 두 사이클 실전송함", BLUE),
    ("2", "7종 센서 14개 CSV 전송", "Accel·Gravity·Gyro·HeartRate·Light·PpgGreen·StepCount × 2사이클 (~6.3MB)", BLUE),
    ("3", "베데스다 담당자 수신 확인", "전송 시각(15:00/15:30)·파일 일치 확인, 받은 14개 파일 역첨부로 회신함", GREEN),
    ("4", "직전 미결 항목(적재 확인) 종료", "서버가 정상 시간·정상 파일로 수신함을 확정 → 전송 검증 단계 종료", GREEN),
    ("5", "운영 안정화 코드 다수 반영", "배터리 실측 가드 · KST 타임스탬프 · DCT 폴더 분리 등 (다음 장)", ACCENT),
]
ty = Inches(2.0)
for num, head, body, color in steps:
    add_rect(s, Inches(0.9), ty, Inches(0.55), Inches(0.55), color)
    add_text(s, Inches(0.9), ty + Inches(0.06), Inches(0.55), Inches(0.45),
             num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.65), ty - Inches(0.02), Inches(10.8), Inches(0.4),
             head, size=17, bold=True, color=NAVY)
    add_text(s, Inches(1.65), ty + Inches(0.34), Inches(10.8), Inches(0.35),
             body, size=13, color=SUB_GRAY)
    ty += Inches(0.96)

# ============ 4. 핵심 성과 (강조) ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(1.0), Inches(0.9), Inches(0.8), Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(1.1), Inches(11.3), Inches(0.6),
         "이번 보고 핵심 성과", size=30, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.9),
         "종단간 자동 전송 + 베데스다 수신 확인 완료",
         size=26, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(2.95), Inches(11.3), Inches(0.5),
         "30분 사이클로 다중 센서 실전송 → 담당자가 정상 시간·정상 파일 수신 확정 (06-17 적재 미결 종료)", size=15, color=RGBColor(0xC8,0xD4,0xE4))
cards = [
    ("2", "전송 사이클 (15:00 / 15:30)", GREEN),
    ("14", "수신 CSV 파일 (7종 × 2)", GREEN),
    ("확인", "담당자 수신 회신 도착", ACCENT),
]
cx = Inches(1.0)
for big, desc, color in cards:
    add_rect(s, cx, Inches(3.9), Inches(3.7), Inches(1.9), RGBColor(0x16, 0x33, 0x55))
    add_rect(s, cx, Inches(3.9), Inches(3.7), Inches(0.12), color)
    add_text(s, cx, Inches(4.35), Inches(3.7), Inches(0.8),
             big, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)
    add_text(s, cx, Inches(5.25), Inches(3.7), Inches(0.4),
             desc, size=13, color=RGBColor(0xC8,0xD4,0xE4), align=PP_ALIGN.CENTER)
    cx += Inches(4.05)
add_text(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.5),
         "→ 06-17의 '단독 1건 + 적재 미확인'을 넘어, 실전송 데이터를 베데스다가 받았음을 회신으로 확정함",
         size=14, bold=True, color=ACCENT)

# ============ 5. 수신 데이터 증빙 ============
s = slide()
header(s, 5, TOTAL, "수신 데이터 증빙", "베데스다가 받은 파일 = 앱이 보낸 파일 — 7종 센서 2사이클 모두 일치")
# 상단 결론 배너
add_rect(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(0.7), RGBColor(0xE3, 0xF1, 0xE8))
add_rect(s, Inches(0.9), Inches(1.85), Inches(0.1), Inches(0.7), GREEN)
add_text(s, Inches(1.2), Inches(2.02), Inches(11.1), Inches(0.4),
         "담당자가 받은 14개 파일을 역첨부 → 전송 시각·센서 종류·파일 수 1:1 일치 확인됨",
         size=15, bold=True, color=NAVY)
# 수신 파일 테이블 (센서 | 15:00 | 15:30)
TX, TY, TW = Inches(0.9), Inches(2.85), Inches(11.6)
COL1, COL2, COL3 = Inches(5.0), Inches(3.3), Inches(3.3)
# 헤더 행
add_rect(s, TX, TY, COL1, Inches(0.5), NAVY)
add_rect(s, TX + COL1, TY, COL2, Inches(0.5), BLUE)
add_rect(s, TX + COL1 + COL2, TY, COL3, Inches(0.5), BLUE)
add_text(s, TX + Inches(0.2), TY + Inches(0.1), COL1, Inches(0.4), "센서 종류", size=14, bold=True, color=WHITE)
add_text(s, TX + COL1, TY + Inches(0.1), COL2, Inches(0.4), "15:00 사이클", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, TX + COL1 + COL2, TY + Inches(0.1), COL3, Inches(0.4), "15:30 사이클", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
sensors = [
    ("Accelerometer", "789 KB", "512 KB"),
    ("Gravity", "779 KB", "464 KB"),
    ("Gyroscope", "658 KB", "393 KB"),
    ("HeartRate", "44 KB", "25 KB"),
    ("Light", "308 KB", "182 KB"),
    ("PpgGreen", "1.27 MB", "753 KB"),
    ("StepCount", "64 KB", "39 KB"),
]
ry = TY + Inches(0.5)
for idx, (name, a, b) in enumerate(sensors):
    bg = WHITE if idx % 2 == 0 else LIGHT_BG
    add_rect(s, TX, ry, COL1, Inches(0.44), bg)
    add_rect(s, TX + COL1, ry, COL2, Inches(0.44), bg)
    add_rect(s, TX + COL1 + COL2, ry, COL3, Inches(0.44), bg)
    add_text(s, TX + Inches(0.2), ry + Inches(0.08), COL1, Inches(0.35), name, size=13, bold=True, color=NAVY, font=FONT_MONO)
    add_text(s, TX + COL1, ry + Inches(0.09), COL2, Inches(0.35), a, size=12, color=GRAY_TEXT, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, TX + COL1 + COL2, ry + Inches(0.09), COL3, Inches(0.35), b, size=12, color=GRAY_TEXT, font=FONT_MONO, align=PP_ALIGN.CENTER)
    ry += Inches(0.44)
# 하단 합계 지표
add_rect(s, TX, ry, TW, Inches(0.55), NAVY)
add_text(s, TX + Inches(0.2), ry + Inches(0.12), TW, Inches(0.4),
         [("합계   7종 센서 · 14개 CSV · 약 6.3 MB", {"size": 14, "bold": True, "color": WHITE})])
add_text(s, Inches(0.9), Inches(6.95), Inches(11.6), Inches(0.35),
         "원본: mails/attachments/2026-06-30_iCReaT_DCT_sensor_upload.zip",
         size=11, color=SUB_GRAY, font=FONT_MONO)

# ============ 6. 06-17 이후 반영한 코드 (운영 안정화) ============
s = slide()
header(s, 6, TOTAL, "06-17 이후 반영한 코드", "실전송 품질을 위한 운영 안정화 수정")
rows = [
    ("배터리 실측 가드", "placeholder 100 차단 — 워치 실측 수신 전 업로드 보류, 음수/100초과 무시", GREEN),
    ("KST 타임스탬프 고정", "로그·파일명 타임스탬프를 KST로 고정 → 전송 시각 정합", BLUE),
    ("DCT 데이터 폴더 분리", "DCT 전송 데이터를 sensor_data_DCT로 분리 → 기존 수집과 충돌 방지", BLUE),
    ("앱 동시설치 분리", "applicationId ...dct 접미사 → 기존 Sensor Monitor와 나란히 설치", BLUE),
    ("무인증 진입 UI 정리", "로그인 화면 이메일 입력칸 제거, \"로그인\"→\"시작\" (DCT 무인증 정책)", BLUE),
    ("셋업 자동화", "setup_phone.ps1 — 설치·권한·배터리 최적화 제외 1회 실행", BLUE),
]
ry = Inches(1.95)
for fname, desc, color in rows:
    add_rect(s, Inches(0.9), ry, Inches(3.6), Inches(0.72), LIGHT_BG)
    add_rect(s, Inches(0.9), ry, Inches(0.1), Inches(0.72), color)
    add_text(s, Inches(1.15), ry + Inches(0.18), Inches(3.3), Inches(0.45),
             fname, size=13, bold=True, color=NAVY)
    add_text(s, Inches(4.7), ry + Inches(0.06), Inches(7.8), Inches(0.6),
             desc, size=12, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    ry += Inches(0.82)

# ============ 7. 06-17 → 07-01 상태 비교 ============
s = slide()
header(s, 7, TOTAL, "06-17 → 07-01 상태 비교", "미결 적재 확인 종료 + 종단간 다중 전송 검증")
add_rect(s, Inches(0.9), Inches(2.0), Inches(3.9), Inches(0.55), NAVY)
add_rect(s, Inches(4.85), Inches(2.0), Inches(3.75), Inches(0.55), RGBColor(0x88,0x88,0x88))
add_rect(s, Inches(8.65), Inches(2.0), Inches(3.85), Inches(0.55), GREEN)
add_text(s, Inches(0.9), Inches(2.1), Inches(3.9), Inches(0.4), "항목", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.85), Inches(2.1), Inches(3.75), Inches(0.4), "06-17", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.65), Inches(2.1), Inches(3.85), Inches(0.4), "07-01", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
comp = [
    ("전송 방식", "curl 단독 1건", "앱 종단간 30분 사이클"),
    ("센서 범위", "단일 샘플", "7종 × 2사이클 14개"),
    ("서버 수신·적재", "미확인 (요청 중)", "담당자 회신 확정"),
    ("타임스탬프/폴더", "기본", "KST 고정 · DCT 폴더 분리"),
]
ry = Inches(2.62)
for item, before, after in comp:
    add_rect(s, Inches(0.9), ry, Inches(3.9), Inches(0.78), LIGHT_BG)
    add_rect(s, Inches(4.85), ry, Inches(3.75), Inches(0.78), RGBColor(0xF0,0xE6,0xE0))
    add_rect(s, Inches(8.65), ry, Inches(3.85), Inches(0.78), RGBColor(0xE3,0xF1,0xE8))
    add_text(s, Inches(1.1), ry + Inches(0.22), Inches(3.6), Inches(0.4), item, size=13, bold=True, color=NAVY)
    add_text(s, Inches(5.0), ry + Inches(0.22), Inches(3.5), Inches(0.4), before, size=12, color=SUB_GRAY)
    add_text(s, Inches(8.8), ry + Inches(0.22), Inches(3.6), Inches(0.4), [("✔ " + after, {"color": GREEN, "bold": True, "size": 12})], size=12)
    ry += Inches(0.86)

# ============ 8. 앞으로 할 것 ============
s = slide()
header(s, 8, TOTAL, "앞으로 할 것", "다음 목표 = QR 스캔 (하드코딩 ID 제거)")
phases = [
    ("다음", "QR 스캔 기능 추가 — 카메라로 Study/Subject 자동 입력, 테스트 하드코딩(C250005/121-001) 제거함. 작년 앱(pbcr_source) QR 매칭 구조 접목", ACCENT),
    ("후속", "인증(세션/API Key) 및 중복·오류코드(401/405) 운영 대응은 업로드 안정화 후 논의함", NAVY),
]
py2 = Inches(2.4)
for tag, body, color in phases:
    add_rect(s, Inches(0.9), py2, Inches(2.2), Inches(1.4), color)
    add_text(s, Inches(0.9), py2 + Inches(0.5), Inches(2.2), Inches(0.5),
             tag, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.25), py2, Inches(9.25), Inches(1.4), LIGHT_BG)
    add_text(s, Inches(3.55), py2 + Inches(0.22), Inches(8.8), Inches(1.0),
             body, size=15, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    py2 += Inches(1.7)
# QR 흐름 한 줄 도식
add_rect(s, Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.85), CODE_BG)
add_text(s, Inches(1.1), Inches(6.32), Inches(11.2), Inches(0.4),
         "QR 스캔 → studyId / subjectId 자동 주입 → 하드코딩 제거 → 다수 대상자 운영 전환",
         size=15, bold=True, color=CODE_FG, font=FONT_MONO, align=PP_ALIGN.CENTER)

# ============ 9. 요약 ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(1.0), Inches(1.0), Inches(0.8), Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.7),
         "요약", size=32, bold=True, color=WHITE)
summary = [
    "06-17엔 curl 단독 1건만 성공했고, 실제 서버 적재 확인이 미결이었음",
    "06-30 워치→앱→서버 30분 사이클로 7종 센서 14개 CSV를 종단간 실전송했음",
    "베데스다 담당자가 정상 시간·정상 파일 수신을 회신 + 받은 파일 압축 역첨부했음",
    "→ 직전 미결(적재 확인) 종료. 종단간 전송 검증 단계 사실상 완료함",
    "다음 목표는 QR 스캔 — Study/Subject 자동 입력으로 테스트 하드코딩 ID 제거함",
]
sy2 = Inches(2.5)
for line in summary:
    add_text(s, Inches(1.0), sy2, Inches(0.3), Inches(0.4), "▸", size=18, bold=True, color=ACCENT)
    add_text(s, Inches(1.4), sy2, Inches(11.0), Inches(0.6), line, size=18, color=WHITE)
    sy2 += Inches(0.78)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides))
