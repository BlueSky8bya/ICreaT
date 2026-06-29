# -*- coding: utf-8 -*-
"""2026-05-27 진행상황 보고 PPT 생성 (발표일 = 오늘 = 27일 기준).
목적: 2026-05-13 보고 이후의 경과 — 무엇이 진행됐고 어떤 코드를 수정했는지.
말투: ~임 / ~했음 (개조식).
스타일: 2026-05-13 덱과 동일 팔레트/폰트. 과도한 기술 세부(인증 방식 등)는 생략.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import os
# 산출물은 reports/decks/에 저장(builders/ 기준 상대경로). 저장소 이동에도 안전. (경로 버그 수정 2026-06-29, 폴더정리 2026-06-29)
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "decks", "[ICreaT과제] 진행상황 보고서_2026-05-27.pptx")

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xE8, 0x74, 0x22)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_TEXT = RGBColor(0x33, 0x33, 0x33)
SUB_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE0, 0xE6, 0xF0)
DIM = RGBColor(0x8A, 0x95, 0xA8)

FONT_KR = "맑은 고딕"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def add_rect(s, x, y, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(s, x, y, w, h, lines, *, font=FONT_KR, size=18, bold=False,
             color=GRAY_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=2):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        # tuple = (text, dict overrides)
        if isinstance(line, tuple):
            text, ov = line
        else:
            text, ov = line, {}
        run = p.add_run()
        run.text = text
        run.font.name = ov.get("font", font)
        run.font.size = Pt(ov.get("size", size))
        run.font.bold = ov.get("bold", bold)
        run.font.color.rgb = ov.get("color", color)
    return tb


def header(s, idx, total, title, subtitle):
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_text(s, Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.72), Inches(0.98), Inches(11.8), Inches(0.35),
             subtitle, size=14, color=SUB_GRAY)
    add_rect(s, Inches(0.7), Inches(1.42), Inches(0.8), Inches(0.06), ACCENT)
    add_text(s, Inches(0.7), Inches(7.08), Inches(8), Inches(0.3),
             "Sensor Monitor → 베데스다 iCReaT DCT 서버 전환", size=10, color=SUB_GRAY)
    add_text(s, Inches(11.4), Inches(7.08), Inches(1.5), Inches(0.3),
             f"{idx} / {total}", size=10, color=SUB_GRAY, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, *, gap=0.52, size=14, dot_color=ACCENT,
            head_color=NAVY, body_color=GRAY_TEXT):
    cy = y
    for head, body in items:
        add_text(s, x, cy, Inches(0.25), Inches(0.3), "▸", size=size, bold=True, color=dot_color)
        add_text(s, x + Inches(0.32), cy - Inches(0.02), w - Inches(0.32), Inches(0.3),
                 head, size=size, bold=True, color=head_color)
        if body:
            add_text(s, x + Inches(0.32), cy + Inches(0.26), w - Inches(0.32), Inches(0.3),
                     body, size=size - 3, color=body_color)
        cy += Inches(gap)


TOTAL = 9

# ============ 1. 타이틀 ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, Inches(5.0), SLIDE_W, Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0),
         "ICReaT 과제 진행상황 보고", size=40, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.6),
         "Sensor Monitor → 베데스다 iCReaT DCT 서버 CSV 업로드 전환", size=18, color=RGBColor(0xC8, 0xD4, 0xE4))
add_text(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.5),
         "2026-05-27  ·  직전 보고(2026-05-13) 이후 경과 중심", size=14, color=RGBColor(0x9F, 0xB0, 0xC8))

# ============ 2. 2주 전 시점 요약 ============
s = slide()
header(s, 2, TOTAL, "2주 전(05-13) 시점", "조사·협의 단계였음 — 코드는 아직 손대기 전")
add_rect(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.06), DIM)
bullets(s, Inches(0.9), Inches(2.1), Inches(11.4), [
    ("기존 앱/서버 CSV 업로드 구조를 코드 기준으로 확인했음", "강원대 Python/Django 수신, 앱은 OkHttp multipart 전송"),
    ("베데스다 회신을 분석해 변경사항을 정리했음", "userID → studyId/subjectId 필요, 무인증(a) 먼저"),
    ("구현 가능한 부분과 베데스다 확인 필요 부분을 분리했음", ""),
    ("산출물만 정리한 상태였음", "회신 메일 초안 · multipart 레퍼런스 · 회신 분석 · 샘플 데이터(sensor_data.zip)"),
], gap=0.78, size=16)
add_rect(s, Inches(0.9), Inches(6.2), Inches(11.4), Inches(0.7), LIGHT_BG)
add_rect(s, Inches(0.9), Inches(6.2), Inches(0.1), Inches(0.7), ACCENT)
add_text(s, Inches(1.2), Inches(6.37), Inches(11), Inches(0.4),
         "→ 이 시점엔 코드 수정 전이었고, 베데스다 회신·확인을 요청하기 직전 상태였음", size=14, bold=True, color=NAVY)

# ============ 3. 그 이후 경과 (타임라인) ============
s = slide()
header(s, 3, TOTAL, "05-13 이후 경과", "회신 → 명세 합의 → QR 분석 → 코드 수정 → 테스트 준비")
steps = [
    ("1", "베데스다에 회신 메일 발송", "전송 항목 확인, 테스트용 ID·신규 서버 주소 요청", BLUE),
    ("2", "답변으로 명세 합의", "엔드포인트 추가 확약, Study/Subject ID 필수", BLUE),
    ("3", "QR 매칭 구조 분석", "작년 앱이 QR로 Study/Subject를 매칭하는 방식 확인", BLUE),
    ("4", "앱 코드 1차 수정 완료", "전송 규격을 베데스다 형식으로 전환 — 이번 보고 핵심", ACCENT),
    ("5", "담당자 오늘 오후 테스트 준비 통보", "우리도 테스트 준비 상태로 정렬했음", GREEN),
]
ty = Inches(2.0)
for num, head, body, color in steps:
    add_rect(s, Inches(0.9), ty, Inches(0.55), Inches(0.55), color)
    add_text(s, Inches(0.9), ty + Inches(0.06), Inches(0.55), Inches(0.45),
             num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.65), ty - Inches(0.02), Inches(10.8), Inches(0.4),
             head, size=17, bold=True, color=NAVY)
    add_text(s, Inches(1.65), ty + Inches(0.34), Inches(10.8), Inches(0.35),
             body, size=13, color=SUB_GRAY)
    ty += Inches(0.96)

# ============ 4. 베데스다 합의 내용 ============
s = slide()
header(s, 4, TOTAL, "베데스다 합의 내용", "내일 테스트의 전제 — 메일로 확정한 사항")
# 좌: 5필드
add_rect(s, Inches(0.9), Inches(2.0), Inches(5.8), Inches(4.6), CODE_BG)
add_text(s, Inches(1.2), Inches(2.2), Inches(5.2), Inches(0.4),
         "전송 multipart 필드 (5종)", size=15, bold=True, color=ACCENT)
fields = [
    ("csvfile", "센서별 30분 병합 CSV"),
    ("studyId", "Study ID (신규)"),
    ("subjectId", "Subject ID (신규)"),
    ("battery", "워치 배터리 잔량"),
    ("timestamp", "전송 시각, Unix epoch 초"),
]
fy = Inches(2.75)
for name, desc in fields:
    add_text(s, Inches(1.2), fy, Inches(2.2), Inches(0.35),
             name, size=15, bold=True, color=CODE_FG, font=FONT_MONO)
    add_text(s, Inches(3.3), fy + Inches(0.02), Inches(3.2), Inches(0.35),
             desc, size=12, color=DIM)
    fy += Inches(0.62)
# 우: 정책 합의
add_text(s, Inches(7.1), Inches(2.0), Inches(5.4), Inches(0.4),
         "정책 합의", size=15, bold=True, color=NAVY)
bullets(s, Inches(7.1), Inches(2.55), Inches(5.4), [
    ("별도 인증 없이 먼저 진행함", "데이터 수집 연동부터 검증"),
    ("기존 userID는 제거함", "여러 과제 운영 → Study/Subject ID 필수"),
    ("전송 성공은 정상 응답으로 판단함", ""),
    ("Study/Subject는 QR 매칭이 최종형", "작년 앱 구조 참조"),
    ("QR 전까지 테스트용 ID로 검증함", "베데스다가 테스트값 제공 예정"),
], gap=0.82, size=14)

# ============ 5. 이번에 수정한 코드 ============
s = slide()
header(s, 5, TOTAL, "이번에 수정한 코드", "QR보다 먼저 — 베데스다 서버가 데이터를 받는지 검증하기 위한 최소 변경, 5개 파일")
rows = [
    ("DeviceInfo.kt", "식별자를 userID → Study/Subject ID로 전환, 테스트용 값 자리 마련"),
    ("ServerConnection.kt", "전송 주소를 신규 베데스다 서버로, 전송 항목을 합의 형식으로 교체"),
    ("AcceptService.kt", "30분 주기 업로드 호출부를 새 식별자 기준으로 교체"),
    ("LoginActivity.kt", "로그인 화면이 새 구조로 바로 진입하도록 변경"),
    ("SensorActivity.kt", "진입 시 Study/Subject 식별자로 초기화"),
]
ry = Inches(2.0)
for fname, desc in rows:
    add_rect(s, Inches(0.9), ry, Inches(3.6), Inches(0.82), LIGHT_BG)
    add_rect(s, Inches(0.9), ry, Inches(0.1), Inches(0.82), NAVY)
    add_text(s, Inches(1.15), ry + Inches(0.2), Inches(3.3), Inches(0.45),
             fname, size=14, bold=True, color=NAVY, font=FONT_MONO)
    add_text(s, Inches(4.7), ry + Inches(0.13), Inches(7.8), Inches(0.6),
             desc, size=14, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    ry += Inches(0.96)

# ============ 6. 전송 페이로드 변경 ============
s = slide()
header(s, 6, TOTAL, "전송 페이로드 변경", "구 강원대(HTTP) → 베데스다(HTTPS), 필드 4종 → 5종")
# before
add_rect(s, Inches(0.9), Inches(2.1), Inches(5.6), Inches(3.6), CODE_BG)
add_text(s, Inches(1.15), Inches(2.3), Inches(5.1), Inches(0.4),
         "변경 전", size=15, bold=True, color=DIM)
add_text(s, Inches(1.15), Inches(2.8), Inches(5.2), Inches(2.7), [
    ("POST http://114.70.120.121:443", {"color": DIM}),
    ("        /forUser/postCurrentData/", {"color": DIM}),
    ("", {}),
    ("csvfile", {"color": CODE_FG}),
    ("userID", {"color": RGBColor(0xE0,0x6C,0x6C)}),
    ("battery", {"color": CODE_FG}),
    ("timestamp", {"color": CODE_FG}),
], size=15, font=FONT_MONO, space_after=4)
# arrow
add_text(s, Inches(6.6), Inches(3.6), Inches(0.7), Inches(0.6),
         "→", size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
# after
add_rect(s, Inches(7.3), Inches(2.1), Inches(5.2), Inches(3.6), CODE_BG)
add_text(s, Inches(7.55), Inches(2.3), Inches(4.7), Inches(0.4),
         "변경 후", size=15, bold=True, color=ACCENT)
add_text(s, Inches(7.55), Inches(2.8), Inches(4.8), Inches(2.7), [
    ("POST https://icreatdct.btsd.io", {"color": GREEN}),
    ("        /<신규 경로>", {"color": GREEN}),
    ("", {}),
    ("csvfile", {"color": CODE_FG}),
    ("studyId", {"color": RGBColor(0x8B,0xD4,0x9B)}),
    ("subjectId", {"color": RGBColor(0x8B,0xD4,0x9B)}),
    ("battery", {"color": CODE_FG}),
    ("timestamp", {"color": CODE_FG}),
], size=15, font=FONT_MONO, space_after=4)
add_text(s, Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.8),
         "식별자 구조도 단일 userID → Study/Subject ID로 바꿨음. "
         "지금은 테스트용 값으로 검증하고, 이후 QR 스캔으로 자동 입력되도록 확장할 예정임.",
         size=13, color=SUB_GRAY)

# ============ 7. 내일 테스트 준비 상태 ============
s = slide()
header(s, 7, TOTAL, "오늘 오후 테스트 준비 상태", "코드는 준비 완료 — 두 가지만 반영하면 바로 테스트 가능")
# 값 2개
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.4),
         "베데스다로부터 받는 즉시 반영할 항목", size=16, bold=True, color=NAVY)
boxes = [
    ("①", "신규 서버 주소", "전송 대상 엔드포인트"),
    ("②", "테스트용 ID", "Study ID / Subject ID 한 쌍"),
]
bx = Inches(0.9)
for num, title, body in boxes:
    add_rect(s, bx, Inches(2.5), Inches(5.7), Inches(1.5), LIGHT_BG)
    add_rect(s, bx, Inches(2.5), Inches(0.1), Inches(1.5), ACCENT)
    add_text(s, bx + Inches(0.3), Inches(2.7), Inches(0.6), Inches(0.6),
             num, size=30, bold=True, color=ACCENT)
    add_text(s, bx + Inches(1.0), Inches(2.72), Inches(4.5), Inches(0.4),
             title, size=16, bold=True, color=NAVY)
    add_text(s, bx + Inches(1.0), Inches(3.22), Inches(4.6), Inches(0.7),
             body, size=13, color=GRAY_TEXT, space_after=1)
    bx += Inches(6.0)
# 절차
add_text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.4),
         "테스트 절차", size=16, bold=True, color=NAVY)
bullets(s, Inches(0.9), Inches(4.95), Inches(11.4), [
    ("받은 주소로 단독 전송 검증함", "준비된 샘플 CSV 사용"),
    ("값 반영 후 디버그 빌드 → 30분 주기 종단간 전송함", ""),
    ("전송 성공 로그 + 베데스다 측 적재 여부 확인함", "담당자 협조"),
], gap=0.62, size=15)

# ============ 8. 앞으로 할 것 ============
s = slide()
header(s, 8, TOTAL, "앞으로 할 것", "오늘 → QR 스캔 → 후속")
phases = [
    ("오늘", "서버 주소·테스트 ID 수신 → 단독 검증 → 앱 종단간 테스트 → 적재 확인함", ACCENT),
    ("다음", "QR 스캔 기능 추가 — 카메라로 QR을 읽어 Study/Subject 자동 입력, 테스트용 값 제거함", BLUE),
    ("후속", "인증 방식 등은 업로드 검증이 안정화된 뒤 논의함", NAVY),
]
py2 = Inches(2.2)
for tag, body, color in phases:
    add_rect(s, Inches(0.9), py2, Inches(2.2), Inches(1.1), color)
    add_text(s, Inches(0.9), py2 + Inches(0.36), Inches(2.2), Inches(0.5),
             tag, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.25), py2, Inches(9.25), Inches(1.1), LIGHT_BG)
    add_text(s, Inches(3.55), py2 + Inches(0.18), Inches(8.8), Inches(0.8),
             body, size=15, color=GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    py2 += Inches(1.4)

# ============ 9. 요약 ============
s = slide()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(1.0), Inches(1.0), Inches(0.8), Inches(0.08), ACCENT)
add_text(s, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.7),
         "요약", size=32, bold=True, color=WHITE)
summary = [
    "2주 전엔 조사·협의 단계였고, 코드는 아직 손대기 전이었음",
    "이후 베데스다와 명세를 합의했고, QR 매칭 구조를 분석했음",
    "앱 코드 1차 수정을 완료했음 — 전송 규격을 베데스다 형식으로 전환, 5개 파일",
    "오늘 테스트는 서버 주소·테스트 ID만 받으면 바로 가능한 상태임",
    "QR 스캔과 인증 등은 업로드 검증 이후 순서로 진행함",
]
sy2 = Inches(2.5)
for line in summary:
    add_text(s, Inches(1.0), sy2, Inches(0.3), Inches(0.4), "▸", size=18, bold=True, color=ACCENT)
    add_text(s, Inches(1.4), sy2, Inches(11.0), Inches(0.6), line, size=18, color=WHITE)
    sy2 += Inches(0.78)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides))
